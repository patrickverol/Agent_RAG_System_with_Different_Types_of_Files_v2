"""
RAG Module
This module implements the Retrieval-Augmented Generation (RAG) system.
It handles document indexing, text extraction from various file formats,
and integration with the vector database for semantic search.
"""

# Import docx module for Word file manipulation
import docx

# Import PyPDF2 module for PDF file manipulation
import PyPDF2

# Import Presentation from pptx package for PowerPoint file manipulation
from pptx import Presentation

# Import pandas for CSV and Excel processing
import pandas as pd

# Import TokenTextSplitter for text tokenization
from langchain_text_splitters import TokenTextSplitter

# Import HuggingFaceEmbeddings for creating embeddings
from langchain_huggingface import HuggingFaceEmbeddings

# Import Qdrant classes for vector database operations
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams

# Import QdrantVectorStore for vector database integration
from langchain_qdrant import QdrantVectorStore

# Import storage module
from storage import get_storage

# Import database auxiliary functions
from aux_functions import get_all_tables_data

# Import chunk ID generation and duplicate checking
from aux_functions import gera_chunk_id, verifica_chunk_existente

# Import required modules
import os
import tempfile


def carrega_texto_word(arquivoname):
    """
    Load text content from a Word document.
    
    Args:
        arquivoname (str): Path to the Word document
        
    Returns:
        str: Extracted text content from the document
    """
    doc = docx.Document(arquivoname)
    fullText = [para.text for para in doc.paragraphs]
    return '\n'.join(fullText)


def carrega_texto_pptx(arquivoname):
    """
    Load text content from a PowerPoint presentation.
    
    Args:
        arquivoname (str): Path to the PowerPoint file
        
    Returns:
        str: Extracted text content from the presentation
    """
    prs = Presentation(arquivoname)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                fullText.append(shape.text)
    return '\n'.join(fullText)


def main_indexing(storage_config):
    """
    Main function for document indexing.
    Processes documents from storage, extracts text, and indexes them in the vector database.
    
    Args:
        storage_config (dict): Configuration for document storage
            - storage_type: Type of storage ('local', 's3', or 'http')
            - base_path: Base path for local storage
            - bucket_name: S3 bucket name (for S3 storage)
            - region_name: AWS region (for S3 storage)
            - endpoint_url: S3 endpoint URL (for S3 storage)
            - base_url: Base URL for HTTP storage
    """
    print(f"\nStorage config: {storage_config}")
    
    # Initialize storage
    storage = get_storage(**storage_config)
    
    # Define model parameters for embeddings
    model_name = "sentence-transformers/msmarco-bert-base-dot-v5"
    model_kwargs = {'device': 'cpu'}
    encode_kwargs = {'normalize_embeddings': True}

    # Initialize HuggingFace embeddings
    hf = HuggingFaceEmbeddings(
        model_name=model_name,
        model_kwargs=model_kwargs,
        encode_kwargs=encode_kwargs
    )

    # Initialize Qdrant client
    client = QdrantClient("http://qdrant:6333")
    collection_name = "RAGVectorDB"

    # Delete collection if it exists
    if client.collection_exists(collection_name):
        print(f"Deleting existing collection: {collection_name}")
        client.delete_collection(collection_name)

    # Create new collection with specified parameters
    print(f"Creating new collection: {collection_name}")
    client.create_collection(
        collection_name,
        vectors_config=VectorParams(size=768, distance=Distance.COSINE)
    )

    # Initialize Qdrant instance
    qdrant = QdrantVectorStore(
        client=client,
        collection_name=collection_name,
        embedding=hf,
        distance=Distance.COSINE
    )

    print("\nIndexing documents...\n")

    # Get list of all documents
    lista_arquivos = storage.list_documents()
    print(f"Found {len(lista_arquivos)} documents to index")
    
    # Process each file in the list
    for arquivo in lista_arquivos:
        print(f"\nProcessing file: {arquivo}")
        
        try:
            # Get document from storage
            temp_file = storage.get_document(arquivo)
            print(f"Retrieved file from storage: {temp_file}")
            
            try:
                arquivo_content = ""
                
                # Process PDF files
                if arquivo.endswith(".pdf"):
                    print("Processing PDF file")
                    reader = PyPDF2.PdfReader(temp_file)
                    for page in reader.pages:
                        arquivo_content += " " + page.extract_text()
                
                # Process text files
                elif arquivo.endswith(".txt"):
                    print("Processing text file")
                    with open(temp_file, 'r') as f:
                        arquivo_content = f.read()
                
                # Process Word documents
                elif arquivo.endswith(".docx"):
                    print("Processing Word document")
                    arquivo_content = carrega_texto_word(temp_file)
                
                # Process PowerPoint presentations
                elif arquivo.endswith(".pptx"):
                    print("Processing PowerPoint presentation")
                    arquivo_content = carrega_texto_pptx(temp_file)
                
                # Process CSV files - each row as separate chunk
                elif arquivo.endswith(".csv"):
                    print("Processing CSV file")
                    try:
                        # Try semicolon separator first for files with empty fields
                        df = pd.read_csv(temp_file, encoding='utf-8', sep=';')
                    except UnicodeDecodeError:
                        df = pd.read_csv(temp_file, encoding='latin-1', sep=';')
                    except Exception as e:
                        print(f"Semicolon separator failed: {e}")
                        # Fallback to comma separator
                        try:
                            df = pd.read_csv(temp_file, encoding='utf-8', sep=',')
                        except UnicodeDecodeError:
                            df = pd.read_csv(temp_file, encoding='latin-1', sep=',')
                    
                    # Convert each row to column-value combinations, replacing empty fields with "0"
                    row_texts = []
                    for index, row in df.iterrows():
                        row_parts = []
                        for col in df.columns:
                            value = row[col]
                            if pd.isna(value) or str(value).strip() == '':
                                row_parts.append(f"{col}: 0")
                            else:
                                row_parts.append(f"{col}: {value}")
                        row_texts.append(", ".join(row_parts))
                    
                    print(f"Split into {len(row_texts)} chunks")
                    
                    # Process each row with chunk ID generation and duplicate checking
                    chunks_added = 0
                    for i, row_text in enumerate(row_texts):
                        metadata = {"path": arquivo, "row": i}
                        
                        # Generate chunk ID
                        chunk_id = gera_chunk_id(row_text, metadata)
                        
                        # Check if chunk already exists
                        if not verifica_chunk_existente(chunk_id, client):
                            # Add chunk_id to metadata
                            metadata["chunk_id"] = chunk_id
                            qdrant.add_texts([row_text], metadatas=[metadata])
                            chunks_added += 1
                        else:
                            print(f"  Chunk {chunk_id} already exists, skipping...")
                    
                    print(f"Successfully indexed {chunks_added} new chunks (skipped {len(row_texts) - chunks_added} duplicates)")
                    continue  # Skip the general text processing
                
                # Process Excel files - each row as separate chunk
                elif arquivo.endswith(".xlsx") or arquivo.endswith(".xls"):
                    print("Processing Excel file")
                    excel_file = pd.ExcelFile(temp_file)
                    all_texts = []
                    
                    for sheet_name in excel_file.sheet_names:
                        df = pd.read_excel(temp_file, sheet_name=sheet_name)
                        # Convert each row to column-value combinations
                        for index, row in df.iterrows():
                            row_parts = [f"{col}: {row[col]}" for col in df.columns if pd.notna(row[col])]
                            all_texts.append(", ".join(row_parts))

                    print(f"Split into {len(all_texts)} chunks")
                    
                    # Process each row with chunk ID generation and duplicate checking
                    chunks_added = 0
                    for i, row_text in enumerate(all_texts):
                        metadata = {"path": arquivo, "row": i}
                        
                        # Generate chunk ID
                        chunk_id = gera_chunk_id(row_text, metadata)
                        
                        # Check if chunk already exists
                        if not verifica_chunk_existente(chunk_id, client):
                            # Add chunk_id to metadata
                            metadata["chunk_id"] = chunk_id
                            qdrant.add_texts([row_text], metadatas=[metadata])
                            chunks_added += 1
                        else:
                            print(f"  Chunk {chunk_id} already exists, skipping...")
                    
                    print(f"Successfully indexed {chunks_added} new chunks (skipped {len(all_texts) - chunks_added} duplicates)")
                    continue  # Skip the general text processing
                
                else:
                    print(f"Skipping unsupported file type: {arquivo}")
                    continue

                # Process PDF, DOCX, PPTX, TXT files - split into 500-token chunks
                print(f"Extracted content length: {len(arquivo_content)} characters")

                # Initialize text splitter
                text_splitter = TokenTextSplitter(chunk_size=500, chunk_overlap=50)
                textos = text_splitter.split_text(arquivo_content)
                print(f"Split into {len(textos)} chunks")
                
                # Process each chunk with chunk ID generation and duplicate checking
                chunks_added = 0
                for i, texto in enumerate(textos):
                    metadata = {"path": arquivo}
                    
                    # Generate chunk ID
                    chunk_id = gera_chunk_id(texto, metadata)
                    
                    # Check if chunk already exists
                    if not verifica_chunk_existente(chunk_id, client):
                        # Add chunk_id to metadata
                        metadata["chunk_id"] = chunk_id
                        qdrant.add_texts([texto], metadatas=[metadata])
                        chunks_added += 1
                    else:
                        print(f"  Chunk {chunk_id} already exists, skipping...")
                
                print(f"Successfully indexed {chunks_added} new chunks (skipped {len(textos) - chunks_added} duplicates)")

            finally:
                # Only delete the temporary file if it's in the temp directory
                if tempfile.gettempdir() in temp_file:
                    os.unlink(temp_file)
                    print(f"Cleaned up temporary file: {temp_file}")

        except Exception as e:
            print(f"Process failed for file {arquivo}: {str(e)}")

    # Process PostgreSQL database data
    print("\nProcessing PostgreSQL database data...")
    try:
        # Get database connection info for metadata
        db_name = os.getenv('POSTGRES_DB', 'rag_db')
        
        # Get all tables data dynamically
        all_tables_data = get_all_tables_data()
        
        if all_tables_data:
            total_rows_indexed = 0
            
            for table_key, table_info in all_tables_data.items():
                schema_name = table_info['schema']
                table_name = table_info['table']
                table_data = table_info['data']
                
                # Create metadata path in format: postgres-database_name-schema_name-table_name
                metadata_path = f"postgres-{db_name}-{schema_name}-{table_name}"
                
                print(f"Processing table: {schema_name}.{table_name} ({len(table_data)} rows)")
                
                # Convert each row to column-value combinations
                row_texts = []
                for record in table_data:
                    row_parts = []
                    for key, value in record.items():
                        if key != 'id':  # Skip id field
                            if value is not None:
                                row_parts.append(f"{key}: {value}")
                            else:
                                row_parts.append(f"{key}: 0")
                    row_texts.append(", ".join(row_parts))

                print(f"Split into {len(row_texts)} chunks from {schema_name}.{table_name}")
                
                # Process each row with chunk ID generation and duplicate checking
                chunks_added = 0
                for i, row_text in enumerate(row_texts):
                    metadata = {
                        "path": metadata_path, 
                        "row": i, 
                        "source": "database", 
                        "schema": schema_name, 
                        "table": table_name
                    }
                    
                    # Generate chunk ID
                    chunk_id = gera_chunk_id(row_text, metadata)
                    
                    # Check if chunk already exists
                    if not verifica_chunk_existente(chunk_id, client):
                        # Add chunk_id to metadata
                        metadata["chunk_id"] = chunk_id
                        qdrant.add_texts([row_text], metadatas=[metadata])
                        chunks_added += 1
                    else:
                        print(f"  Chunk {chunk_id} already exists, skipping...")
                
                total_rows_indexed += chunks_added
                print(f"Successfully indexed {chunks_added} new chunks from {schema_name}.{table_name} (skipped {len(row_texts) - chunks_added} duplicates)")
            
            print(f"Total database rows indexed: {total_rows_indexed}")
        else:
            print("No tables found in database")
            
    except Exception as e:
        print(f"Error processing PostgreSQL data: {e}")

    print("\nIndexing Completed!\n")


if __name__ == "__main__":

    # Get storage type from environment variable
    storage_type = os.getenv('STORAGE_TYPE', 'http')
    
    # Configure storage based on type
    if storage_type == 's3':
        storage_config = {
            'storage_type': 's3',
            'bucket_name': os.getenv('S3_BUCKET_NAME'),
            'region_name': os.getenv('AWS_REGION'),
            'endpoint_url': os.getenv('S3_ENDPOINT_URL')
        }
    elif storage_type == 'local':
        storage_config = {
            'storage_type': 'local',
            'base_path': os.getenv('DOCUMENTS_PATH', '/app/documents')
        }
    else:  # http
        storage_config = {
            'storage_type': 'http',
            'base_url': os.getenv('DOCUMENT_STORAGE_URL', 'http://document_storage:8080')
        }
    
    main_indexing(storage_config)